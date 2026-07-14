#!/usr/bin/env python3
"""Backup slides are not part of the talk, so they are not part of its clock.

Nearly every conference deck carries a backup section — the Q&A slides you do not present
and open only if someone asks. The budget check counted them against the clock, so a
10-minute talk with four backup slides was told to cut, and the honest way to satisfy the
tool was to delete the Q&A preparation. That is the tool being wrong in the one place a
speaker most needs to be prepared.

The clock now stops at the first slide marked "Backup" / "Appendix" / "Q&A" / "백업".
Density and type size still apply past that line — a backup slide is shown *under
questioning*, which is the worst possible moment to find out it is a wall of 11-pt text.

Skips cleanly (exit 0) without python-pptx. Network-free.

    python3 skills/present-paper/tests/test_backup_off_the_clock.py
"""
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx not installed — SKIP (compile-only)")
    sys.exit(0)

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))
import check_deck_budget as budget  # noqa: E402

fails: list[str] = []


def deck(headlines, body_words=8, pt=24):
    """One slide per headline: the headline, plus a little body so it counts as content."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    for h in headlines:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        for txt, y in ((h, 0.5), (" ".join(["word"] * body_words), 2.5)):
            if not txt:
                continue
            tb = s.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11), Inches(1.2))
            tb.text_frame.text = txt
            tb.text_frame.paragraphs[0].runs[0].font.size = Pt(pt)
    p = Path(tempfile.mkdtemp()) / "d.pptx"
    prs.save(str(p))
    return p


def verdicts(path, minutes=10, archetype="conference_oral"):
    return [f.verdict for f in budget.audit(path, archetype, minutes)]


# --- the boundary itself ------------------------------------------------------------
CASES = [
    (["Backup"], 0),
    (["Backup — Q&A"], 0),
    (["Appendix"], 0),
    (["Q & A"], 0),
    (["백업"], 0),
    # a sentence that merely contains the word is a sentence, not a signpost
    (["Reserved for the appendix of the guideline"], None),
    (["Supplemental oxygen was given to every patient"], None),
    (["58 tells us nothing — four numbers do"], None),
]
for heads, want in CASES:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    for h in heads:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(1))
        tb.text_frame.text = h
    p = Path(tempfile.mkdtemp()) / "b.pptx"
    prs.save(str(p))
    slides, _n, _w, _h = budget.read_deck(p)
    got = budget.find_backup_boundary(slides)
    if got != want:
        fails.append(f"boundary for {heads[0]!r}: want {want}, got {got}")


# --- the clock stops there ----------------------------------------------------------
TALK = [f"Finding {i} changed what we do" for i in range(11)]      # 11 presented slides
BACKUP = ["Backup"] + [f"Answer to question {i}" for i in range(6)]  # 7 more, off the clock

# 11 presented slides for a 10-min conference oral: allowed 10, slack to 12.5 -> fits.
if "DECK_OVER_BUDGET" in verdicts(deck(TALK)):
    fails.append("11 presented slides for a 10-min oral should fit (budget 10, +25% slack)")

# The same talk with a backup section must STILL fit -- that is the whole fix.
if "DECK_OVER_BUDGET" in verdicts(deck(TALK + BACKUP)):
    fails.append("backup slides were counted against the clock — the bug this test exists for")

# ...but a genuinely over-long talk must still be caught, backup section or not.
if "DECK_OVER_BUDGET" not in verdicts(deck([f"Point {i} matters here" for i in range(20)] + BACKUP)):
    fails.append("20 presented slides for a 10-min oral must still be caught")


# --- legibility does NOT stop there -------------------------------------------------
# A dense, tiny backup slide is shown under questioning. It still has to be readable.
dense = deck(TALK + ["Backup"] + ["Answer"], body_words=0)
prs = Presentation(str(dense))
s = prs.slides[-1]
tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11), Inches(4))
tb.text_frame.word_wrap = True
tb.text_frame.text = " ".join(["word"] * 90)              # over the 60-word ceiling
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(9)     # under the 20 pt floor
prs.save(str(dense))
v = verdicts(dense)
if "SLIDE_TOO_DENSE" not in v:
    fails.append("a 90-word backup slide escaped the density check — backups get shown too")
if "TYPE_TOO_SMALL" not in v:
    fails.append("a 9-pt backup slide escaped the type floor — backups get shown too")
if "DECK_OVER_BUDGET" in v:
    fails.append("the backup section was still on the clock")


if fails:
    print("FAIL — backup section")
    for f in fails:
        print("  ·", f)
    sys.exit(1)
print(f"PASS — boundary ({len(CASES)} cases), clock stops at backup, legibility does not")
