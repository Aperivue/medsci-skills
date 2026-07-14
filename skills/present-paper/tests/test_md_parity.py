#!/usr/bin/env python3
"""The inline-markdown parser must not eat an asterisk that belongs to the content.

Two call sites parse ``**bold**`` / ``*italic*`` — ``scripts/inject_speaker_notes.py``
(speaker notes) and ``templates/build_pptx_nature_lancet.py`` (slide body). SKILL.md
mandates a word-boundary pattern for both. They had drifted to a naive
``\\*[^*\\n]+\\*``, which silently deleted the asterisk in ``DRB1*07:01`` — rewriting a
genotype in a medical deck and losing the bold with it.

This test pins two things:

  1. **Parity** — both modules compile the same pattern, so they cannot drift apart again.
  2. **Behaviour** — an allele survives; nothing is ever dropped.

Skips cleanly (exit 0) without python-pptx. Network-free.

    python3 skills/present-paper/tests/test_md_parity.py
"""
import importlib.util
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation  # noqa: F401
except ImportError:
    print("python-pptx not installed — SKIP (compile-only)")
    sys.exit(0)

ROOT = Path(__file__).resolve().parents[1]


def load(path: Path, name: str):
    spec = importlib.util.spec_from_file_location(name, path)
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    return m


notes_mod = load(ROOT / "scripts" / "inject_speaker_notes.py", "isn")
build_mod = load(ROOT / "templates" / "build_pptx_nature_lancet.py", "bnl")

fails: list[str] = []


# --- 1. parity ---------------------------------------------------------------------
if notes_mod._MD_INLINE.pattern != build_mod._MD_PATTERN.pattern:
    fails.append(
        "the two call sites use DIFFERENT inline-markdown patterns — they will drift:\n"
        f"    inject_speaker_notes : {notes_mod._MD_INLINE.pattern!r}\n"
        f"    build_pptx_nature_l. : {build_mod._MD_PATTERN.pattern!r}"
    )

# The naive pattern is the bug. Refuse it by name, in either site.
NAIVE = r"(\*\*[^*\n]+\*\*|\*[^*\n]+\*)"
for label, pat in (("inject_speaker_notes", notes_mod._MD_INLINE.pattern),
                   ("build_pptx_nature_lancet", build_mod._MD_PATTERN.pattern)):
    if pat == NAIVE:
        fails.append(f"{label} regressed to the naive pattern — it eats DRB1*07:01")
    if "(?![A-Za-z0-9])" not in pat or "(?<![A-Za-z0-9])" not in pat:
        fails.append(f"{label} italic rule has no word boundary — an allele asterisk will be eaten")


# --- 2. behaviour ------------------------------------------------------------------
# (source, visible text after rendering, spans that must be bold)
CASES = [
    ("Carriers of **DRB1*07:01** had higher risk",
     "Carriers of DRB1*07:01 had higher risk", ["DRB1*07:01"]),
    ("The allele DRB1*04:02 and HLA-A*02:01 were typed",
     "The allele DRB1*04:02 and HLA-A*02:01 were typed", []),
    ("rs1801133*T carriers",
     "rs1801133*T carriers", []),
    ("**bold** and *ital* both",
     "bold and ital both", ["bold"]),
    # nested emphasis: markers stay visible (author's error, made loud) — but the
    # bold span still applies and NOT ONE CHARACTER is lost. The old parser dropped
    # the final "e" of "core" here.
    ("**The sign is *itself* fluid** — a core",
     "The sign is *itself* fluid — a core", ["The sign is *itself* fluid"]),
]


def render(mod_pattern, text):
    """Reproduce the split/emit contract both modules share."""
    out, bolds = [], []
    for part in mod_pattern.split(text):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**") and len(part) > 4:
            out.append(part[2:-2]); bolds.append(part[2:-2])
        elif (part.startswith("*") and part.endswith("*")
              and not part.startswith("**") and len(part) > 2):
            out.append(part[1:-1])
        else:
            out.append(part)
    return "".join(out), bolds


for src, want_text, want_bold in CASES:
    got_text, got_bold = render(notes_mod._MD_INLINE, src)
    if got_text != want_text:
        fails.append(f"rendered text differs\n    in   {src!r}\n    want {want_text!r}\n    got  {got_text!r}")
    if got_bold != want_bold:
        fails.append(f"bold spans differ for {src!r}\n    want {want_bold!r}\n    got  {got_bold!r}")
    # nothing may vanish: every non-marker character survives
    if len(got_text) < len(re.sub(r"\*\*", "", src)) - src.count("*"):
        fails.append(f"characters were dropped rendering {src!r} -> {got_text!r}")


# --- 3. the real notes injector, end to end ----------------------------------------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
tf = slide.notes_slide.notes_text_frame
tf.clear()
notes_mod._add_markdown_line(tf.paragraphs[0], "Carriers of **DRB1*07:01** had higher risk")
runs = tf.paragraphs[0].runs
visible = "".join(r.text for r in runs)
if "DRB1*07:01" not in visible:
    fails.append(f"inject_speaker_notes corrupted the allele in a real note: {visible!r}")
if not any(r.font.bold and "DRB1*07:01" in r.text for r in runs):
    fails.append(f"the allele lost its bold styling: {[(r.text, r.font.bold) for r in runs]!r}")


if fails:
    print("FAIL — inline-markdown parser")
    for f in fails:
        print("  ·", f)
    sys.exit(1)
print(f"PASS — parity + {len(CASES)} rendering cases; alleles survive, nothing dropped")
