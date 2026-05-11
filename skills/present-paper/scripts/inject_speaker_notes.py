#!/usr/bin/env python3
"""Inject speaker notes into PowerPoint presentation slides.

This script adds or replaces speaker notes in a PPTX file without modifying
slide content, layout, or design. Notes are defined as a dictionary mapping
slide numbers (1-indexed) to note text.

Markdown handling (since 2026-05-11):
    python-pptx's ``notes_text_frame.text = ...`` writes the value as *plain
    text* — Markdown syntax like ``**bold**`` and ``*italic*`` is therefore
    rendered verbatim in PowerPoint's Presenter View (raw ``*`` symbols
    visible to the speaker). To avoid that, this script parses Markdown
    inline emphasis and converts it to run-level bold/italic styling.

    Use ``--no-markdown`` to disable parsing (legacy behavior — write as
    raw text). Use ``--cjk-font NAME`` to enforce an East-Asia font on
    every run (default: Apple SD Gothic Neo, required for Korean notes on
    Mac PowerPoint).

Usage:
    python inject_speaker_notes.py input.pptx
    python inject_speaker_notes.py input.pptx -o output.pptx
    python inject_speaker_notes.py input.pptx --append
    python inject_speaker_notes.py input.pptx --dry-run
    python inject_speaker_notes.py input.pptx --no-markdown
    python inject_speaker_notes.py input.pptx --cjk-font "Malgun Gothic"

Requirements:
    pip install python-pptx

License: MIT
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Pt
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


# ---------------------------------------------------------------------------
# Speaker notes dictionary
# Map slide number (1-indexed) to note text.
# Empty string or missing key = skip that slide.
# Inline markdown supported: **bold**, *italic*  (single-line, no nesting).
# ---------------------------------------------------------------------------
notes: dict[int, str] = {
    # 1: """Speaker note for slide 1.""",
    # 2: """Speaker note for slide 2 with **emphasis** and *highlight*.""",
}


# Inline markdown emphasis: **bold** or *italic* (non-greedy, no nesting,
# no line-wrap). `**` is checked first so it is not consumed by single `*`.
_MD_INLINE = re.compile(r"(\*\*[^*\n]+\*\*|\*[^*\n]+\*)")


def _set_eastasia_font(run, font_name: str) -> None:
    """Force the East-Asia (CJK) font on a single run.

    PowerPoint Mac falls back to a system default for Korean glyphs when the
    eastAsia attribute is missing, which often differs from the Latin font
    and breaks visual consistency in Presenter View.
    """
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = rPr.makeelement(qn("a:ea"), {"typeface": font_name})
    rPr.append(ea)


def _add_styled_run(p, text: str, *, bold: bool = False, italic: bool = False,
                    size_pt: int = 12, cjk_font: str | None) -> None:
    """Append a styled run to paragraph ``p``."""
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    if cjk_font:
        _set_eastasia_font(run, cjk_font)


def _fill_notes_with_markdown(tf, text: str, *, size_pt: int = 12,
                              cjk_font: str | None) -> None:
    """Replace the notes text frame contents with parsed markdown.

    Each input line maps to one paragraph. Within a line, ``**bold**`` and
    ``*italic*`` segments become separate runs with the appropriate styling.
    """
    tf.clear()
    first_paragraph = True
    for line in text.split("\n"):
        if first_paragraph:
            p = tf.paragraphs[0]
            first_paragraph = False
        else:
            p = tf.add_paragraph()
        if not line:
            continue
        for part in _MD_INLINE.split(line):
            if not part:
                continue
            if part.startswith("**") and part.endswith("**") and len(part) > 4:
                _add_styled_run(p, part[2:-2], bold=True, size_pt=size_pt,
                                cjk_font=cjk_font)
            elif (part.startswith("*") and part.endswith("*")
                  and len(part) > 2 and not part.startswith("**")):
                _add_styled_run(p, part[1:-1], italic=True, size_pt=size_pt,
                                cjk_font=cjk_font)
            else:
                _add_styled_run(p, part, size_pt=size_pt, cjk_font=cjk_font)


def inject_notes(
    input_path: str,
    output_path: str | None = None,
    append: bool = False,
    dry_run: bool = False,
    markdown: bool = True,
    cjk_font: str | None = "Apple SD Gothic Neo",
) -> None:
    """Inject speaker notes into a PPTX file.

    Args:
        input_path: Path to input PPTX file.
        output_path: Path to output PPTX file. Defaults to input with _notes suffix.
        append: If True, append to existing notes instead of replacing.
        dry_run: If True, print what would be done without saving.
        markdown: If True, parse **bold** / *italic* into run-level styling.
            If False, write the value as plain text (legacy behavior).
        cjk_font: East-Asia font enforced per run when markdown=True.
            Set to None to leave the eastAsia attribute unset.
    """
    input_file = Path(input_path)
    if not input_file.exists():
        print(f"Error: {input_file} not found")
        sys.exit(1)

    if output_path is None:
        output_file = input_file.with_stem(input_file.stem + "_notes")
    else:
        output_file = Path(output_path)

    prs = Presentation(str(input_file))
    total_slides = len(prs.slides)
    updated = 0

    for i, slide in enumerate(prs.slides, 1):
        if i not in notes or not notes[i]:
            continue

        if dry_run:
            preview = notes[i][:80].replace("\n", " ")
            mode = "append" if append else "set"
            mark = "+md" if markdown else "plain"
            print(f"  Slide {i:2d}: would {mode} ({mark}) → {preview}...")
            updated += 1
            continue

        if not slide.has_notes_slide:
            slide.notes_slide  # creates notes slide

        tf = slide.notes_slide.notes_text_frame
        if markdown:
            existing = tf.text if append else ""
            target_text = (existing + "\n\n---\n\n" + notes[i]) if existing.strip() else notes[i]
            _fill_notes_with_markdown(tf, target_text, cjk_font=cjk_font)
        else:
            if append and tf.text.strip():
                tf.text = tf.text + "\n\n---\n\n" + notes[i]
            else:
                tf.text = notes[i]
        updated += 1

    if dry_run:
        print(f"\nDry run: {updated}/{total_slides} slides would be updated")
        return

    prs.save(str(output_file))
    print(f"Done: {output_file} ({updated}/{total_slides} slides updated)")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Inject speaker notes into PowerPoint slides",
        epilog="Notes are defined in the 'notes' dictionary in this script.",
    )
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument(
        "-o", "--output",
        help="Output PPTX file (default: input with _notes suffix)",
    )
    parser.add_argument(
        "--append",
        action="store_true",
        help="Append to existing notes instead of replacing",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Print what would be done without saving",
    )
    parser.add_argument(
        "--no-markdown",
        action="store_true",
        help="Disable Markdown parsing — write notes as plain text (legacy mode).",
    )
    parser.add_argument(
        "--cjk-font",
        default="Apple SD Gothic Neo",
        help="East-Asia font enforced on every run (default: Apple SD Gothic Neo). "
             "Set to '' to leave eastAsia attribute unset.",
    )
    args = parser.parse_args()

    if not notes:
        print("Warning: notes dictionary is empty. Edit this script to add notes.")
        print("Example:")
        print('  notes = {')
        print('      1: """Your note for slide 1.""",')
        print('      2: """Note with **emphasis** and *highlight*.""",')
        print('  }')
        sys.exit(0)

    inject_notes(
        args.input,
        args.output,
        args.append,
        args.dry_run,
        markdown=not args.no_markdown,
        cjk_font=(args.cjk_font or None),
    )


if __name__ == "__main__":
    main()
