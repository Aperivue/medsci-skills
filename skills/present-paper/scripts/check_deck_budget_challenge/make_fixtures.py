#!/usr/bin/env python3
"""Four decks, to test one claim: the same slides are right for one room and wrong for another.

  academic.pptx    ordinary academic slides — ~40 words, 20 pt. Fits a 10-minute oral. As a KEYNOTE
                   the same deck is a wall of text, and that is the point: no single global
                   words-per-slide number could be right for both.
  keynote.pptx     six words and a big number. Would be an empty academic slide; is a good keynote.
  bloated.pptx     60 slides for a 10-minute talk — a talk whose ending gets taken away at the mic.
  tiny_type.pptx   12-pt body text. The back row exists.

Written wherever the caller says (a temp dir, in practice). Nothing built lands in the repo tree.
Needs python-pptx (CI installs it).
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

BLANK = 6

# ~30 words of evidence under a 6-word assertion: an ordinary, well-made academic slide. Note how
# little that actually is — the 40-word working figure is a target, not a licence to fill the space.
ACADEMIC_BODY = (
    "412 consecutive patients. Recurrence 12% versus 26% over a median 3.2 years. "
    "The difference held after adjustment for tumour size and did not vary by centre."
)


def deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs


def box(slide, text, x, y, w, h, pt):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.paragraphs[0].runs[0].font.size = Pt(pt)
    return tb


def build_academic(path: Path) -> None:
    """~40 words a slide at 20 pt: unremarkable for an oral abstract, fatal for a keynote."""
    prs = deck()
    for i in range(9):
        s = prs.slides.add_slide(prs.slide_layouts[BLANK])
        box(s, f"Adjunctive ablation halved local recurrence ({i + 1})", 0.8, 0.8, 11.5, 1.1, 32)
        box(s, ACADEMIC_BODY, 0.8, 2.3, 11.5, 3.0, 20)
        box(s, str(i + 1), 12.6, 7.0, 0.4, 0.3, 10)
    prs.save(str(path))


def build_keynote(path: Path) -> None:
    """The slide is punctuation. A room that came to be moved is not reading."""
    prs = deck()
    for text in ("12% versus 26%", "The tract was the route", "So: ablate the tract",
                 "One patient in eight", "That is the whole talk", "Questions"):
        s = prs.slides.add_slide(prs.slide_layouts[BLANK])
        box(s, text, 1.0, 2.6, 11.0, 2.2, 54)
    prs.save(str(path))


def build_bloated(path: Path) -> None:
    """60 slides for ten minutes. Not a style choice — an ending taken away at the microphone."""
    prs = deck()
    for i in range(60):
        s = prs.slides.add_slide(prs.slide_layouts[BLANK])
        box(s, f"Point number {i + 1} about the cohort", 0.8, 0.8, 11.5, 1.1, 30)
        box(s, "A short supporting line.", 0.8, 2.3, 11.5, 0.8, 22)
    prs.save(str(path))


def build_tiny_type(path: Path) -> None:
    """Below the floor, text is not read. It is seen. That is another word for decoration."""
    prs = deck()
    for i in range(6):
        s = prs.slides.add_slide(prs.slide_layouts[BLANK])
        box(s, f"Cohort characteristics ({i + 1})", 0.8, 0.8, 11.5, 1.0, 30)
        box(s, "Median age 61 years; 54% male; 38% had prior therapy.", 0.8, 2.2, 11.5, 0.8, 12)
    prs.save(str(path))


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.cwd()
    out.mkdir(parents=True, exist_ok=True)
    build_academic(out / "academic.pptx")
    build_keynote(out / "keynote.pptx")
    build_bloated(out / "bloated.pptx")
    build_tiny_type(out / "tiny_type.pptx")
    print(f"wrote 4 decks into {out}")
