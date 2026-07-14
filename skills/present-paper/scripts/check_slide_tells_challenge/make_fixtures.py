#!/usr/bin/env python3
"""Build the two fixture decks the challenge runs against.

They are BUILT, not committed: a .pptx is a binary, and a binary fixture is a fixture nobody ever
reads again. Building it means the tells are written down in Python, where the next person can see
exactly what "an AI-made deck" means here and disagree with me.

  tells.pptx   every mark check_slide_tells looks for, planted on purpose.
  clean.pptx   the same content, made the way a person makes it — and it must come back CLEAN.
               This is the half that matters. A checker that only ever fires is a checker that will
               be turned off.

Needs python-pptx (CI installs it). Writes into the directory given as argv[1].
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent

BLANK = 6  # the blank layout in the default template


def textbox(slide, text, x, y, w, h, pt):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.text = text
    tf.paragraphs[0].runs[0].font.size = Pt(pt)
    return tb


def build_tells(path: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    for i in range(6):
        s = prs.slides.add_slide(prs.slide_layouts[BLANK])

        # CHROME_ON_EVERY_SLIDE — the little words along both edges, on all six slides.
        textbox(s, "NEUROGENETICS · SECTION 2", 0.7, 0.32, 6.0, 0.4, 11)
        textbox(s, "2026 · NEUROGENETICS", 0.7, 7.05, 4.0, 0.35, 9)

        if i == 0:
            # TOPIC_TITLE — a filing label on a slide that has a body.
            textbox(s, "Results", 0.7, 1.0, 8.0, 0.9, 32)
            textbox(s, "Recurrence was lower in the adjunctive arm.", 0.7, 2.2, 8.0, 0.6, 18)
            textbox(s, "The difference persisted at 24 months.", 0.7, 3.0, 8.0, 0.6, 18)

        elif i == 1:
            # SCAFFOLD_PHRASE — the deck narrating its own construction.
            textbox(s, "Adjunctive ablation halved recurrence", 0.7, 1.0, 8.0, 0.9, 32)
            textbox(s, "요약하자면, 국소 재발률이 절반으로 줄었다.", 0.7, 2.2, 9.0, 0.6, 18)
            textbox(s, "The key takeaway is that margins matter.", 0.7, 3.0, 9.0, 0.6, 18)

        elif i == 2:
            # SHAPE_MONOTONY — the same rounded box, over and over.
            textbox(s, "Pipeline", 0.7, 1.0, 8.0, 0.9, 32)
            for k in range(9):
                s.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.8 + (k % 3) * 2.6), Inches(2.2 + (k // 3) * 1.3),
                    Inches(2.2), Inches(0.9),
                )

        elif i == 3:
            # DEAD_SPACE_BAND — a wide empty stripe between the objects.
            textbox(s, "Cohort", 0.7, 0.9, 8.0, 0.8, 30)
            textbox(s, "n = 412", 0.7, 5.9, 3.0, 0.5, 18)
            textbox(s, "median follow-up 3.2 y", 4.5, 5.9, 4.0, 0.5, 18)

        elif i == 4:
            # ARROW_NO_SEMANTICS — three arrows, none of them saying what it means.
            textbox(s, "Mechanism", 0.7, 0.9, 8.0, 0.8, 30)
            for k in range(3):
                s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8 + k * 4.0), Inches(3.0),
                                   Inches(2.4), Inches(1.0))
            for k in range(3):
                s.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(3.3 + k * 4.0), Inches(3.5), Inches(4.7 + k * 4.0), Inches(3.5),
                )

        else:
            textbox(s, "Discussion", 0.7, 1.0, 8.0, 0.9, 32)  # another filing label
            textbox(s, "이는 세 가지 층위에서 살펴볼 수 있다.", 0.7, 2.2, 9.0, 0.6, 18)
            textbox(s, "It is important to note the small sample.", 0.7, 3.0, 9.0, 0.6, 18)

    prs.save(str(path))


def build_clean(path: Path) -> None:
    """The same talk, made by someone who is trying to be understood.

    No edge chrome. Headlines that state findings. Arrows that say what they mean. Shapes that
    differ because the ideas differ. Nothing here should fire — if it does, the checker is wrong,
    not the deck.
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1 — a headline that is the finding, and one figure-sized block of evidence.
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    textbox(s, "Adjunctive ablation halved local recurrence (12% vs 26%)", 0.8, 0.8, 11.0, 1.2, 34)
    textbox(s, "412 patients, median follow-up 3.2 years. The difference held at 24 months and did "
               "not depend on tumour size.", 0.8, 2.4, 11.0, 2.2, 20)
    textbox(s, "1", 12.6, 7.0, 0.4, 0.3, 10)  # a bare page number earns its place

    # 2 — a divider. A filing label IS the job here, and must not be flagged.
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    textbox(s, "Methods", 0.8, 3.0, 8.0, 1.4, 44)

    # 3 — a diagram whose arrows each carry their claim.
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    textbox(s, "Seeding follows the catheter tract, not the pleural space", 0.8, 0.7, 11.0, 1.0, 30)
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.0), Inches(2.6), Inches(1.1))
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(2.9), Inches(2.8), Inches(1.3))
    s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(9.6), Inches(3.05), Inches(2.6), Inches(1.0))
    for k, (x0, x1) in enumerate(((3.6, 5.1), (8.1, 9.5))):
        s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(3.55),
                               Inches(x1), Inches(3.55))
        textbox(s, "seeds along", x0 + 0.05, 3.05, 1.4, 0.35, 12)  # the label sits on the arrow
    textbox(s, "2", 12.6, 7.0, 0.4, 0.3, 10)

    prs.save(str(path))


if __name__ == "__main__":
    # The decks are written wherever the caller says — a temp dir, in practice. Nothing built lands
    # in the repository tree: a .pptx is a binary, it would be inventoried into the release manifest
    # while being (rightly) gitignored, and the classroom ZIP would then fail its own hash check.
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.cwd()
    out.mkdir(parents=True, exist_ok=True)
    build_tells(out / "tells.pptx")
    build_clean(out / "clean.pptx")
    print(f"wrote {out/'tells.pptx'} and {out/'clean.pptx'}")
