#!/usr/bin/env python3
"""Regression test for scripts/validate_pptx_mac_compat.py.

Builds a clean .pptx (python-pptx, with a corrected docProps/app.xml slide
count) that must PASS, then injects each of the four Mac-incompatibility defect
classes into a copy and asserts the validator FAILs under --strict:

    TIFF media, <a:sp3d> 3-D bevel, app.xml slide-count mismatch, srcRect
    over-crop (> 100000).

Also asserts a missing input exits 2. Requires python-pptx (a CI dependency).
"""

from __future__ import annotations

import re
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
VALIDATOR = HERE.parent / "scripts" / "validate_pptx_mac_compat.py"

_pass = 0
_fail = 0


def ck(label: str, expected: int, actual: int) -> None:
    global _pass, _fail
    if expected == actual:
        print(f"  PASS  {label:<48} exit={actual}")
        _pass += 1
    else:
        print(f"  FAIL  {label:<48} expected={expected} actual={actual}")
        _fail += 1


def run(pptx: Path, strict: bool = True) -> int:
    cmd = [sys.executable, str(VALIDATOR), str(pptx)]
    if strict:
        cmd.append("--strict")
    return subprocess.run(cmd, capture_output=True, text=True).returncode


def build_clean(path: Path, n_slides: int = 2) -> None:
    """A python-pptx deck with app.xml <Slides> corrected to match — should PASS."""
    from pptx import Presentation

    prs = Presentation()
    for _ in range(n_slides):
        prs.slides.add_slide(prs.slide_layouts[6])  # blank
    prs.save(str(path))
    # python-pptx writes <Slides>0</Slides>; fix it so the clean baseline passes.
    _rewrite_member(path, "docProps/app.xml",
                     lambda b: re.sub(rb"<Slides>\d+</Slides>",
                                      f"<Slides>{n_slides}</Slides>".encode(), b))


def _read_member(path: Path, member: str) -> bytes:
    with zipfile.ZipFile(path) as z:
        return z.read(member)


def _rewrite_member(path: Path, member: str, transform, extra: dict | None = None) -> None:
    """Rewrite the zip, replacing `member` (via transform) and adding `extra` files."""
    with zipfile.ZipFile(path) as z:
        items = {n: z.read(n) for n in z.namelist()}
    if member in items:
        items[member] = transform(items[member])
    if extra:
        items.update(extra)
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as z:
        for n, data in items.items():
            z.writestr(n, data)
    tmp.replace(path)


def first_slide(path: Path) -> str:
    with zipfile.ZipFile(path) as z:
        for n in sorted(z.namelist()):
            if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
                return n
    raise AssertionError("no slide XML in fixture")


def main() -> int:
    if VALIDATOR.exists() is False:
        print(f"validator missing: {VALIDATOR}")
        return 1
    tmp = Path(tempfile.mkdtemp())

    clean = tmp / "clean.pptx"
    build_clean(clean)
    ck("clean deck (app.xml fixed) passes", 0, run(clean))

    # 1) TIFF embedded in ppt/media/
    tiff = tmp / "tiff.pptx"
    build_clean(tiff)
    _rewrite_member(tiff, "docProps/app.xml", lambda b: b,
                    extra={"ppt/media/image1.tiff": b"II*\x00 fake tiff"})
    ck("TIFF media -> FAIL", 1, run(tiff))

    # 2) <a:sp3d> 3-D bevel inside a slide
    sp3d = tmp / "sp3d.pptx"
    build_clean(sp3d)
    sl = first_slide(sp3d)
    _rewrite_member(sp3d, sl, lambda b: b.replace(b"</p:sld>", b"<a:sp3d/></p:sld>"))
    ck("sp3d bevel -> FAIL", 1, run(sp3d))

    # 3) app.xml slide-count mismatch (declare a wrong count)
    appx = tmp / "appx.pptx"
    build_clean(appx)
    _rewrite_member(appx, "docProps/app.xml",
                    lambda b: re.sub(rb"<Slides>\d+</Slides>", b"<Slides>99</Slides>", b))
    ck("app.xml slide-count mismatch -> FAIL", 1, run(appx))

    # 4) srcRect over-crop (> 100000)
    src = tmp / "srcrect.pptx"
    build_clean(src)
    sl2 = first_slide(src)
    _rewrite_member(src, sl2, lambda b: b.replace(b"</p:sld>", b'<a:srcRect l="997171"/></p:sld>'))
    ck("srcRect over-crop (>100000) -> FAIL", 1, run(src))

    # 5) missing input -> exit 2
    ck("missing input -> exit 2", 2, run(tmp / "does_not_exist.pptx"))

    # 6) non-strict tolerates a WARN-only deck (no app.xml -> WARN, exit 0)
    nowarn = tmp / "noappxml.pptx"
    build_clean(nowarn)
    _rewrite_member(nowarn, "docProps/app.xml", lambda b: b, extra=None)
    # remove app.xml entirely
    with zipfile.ZipFile(nowarn) as z:
        items = {n: z.read(n) for n in z.namelist() if n != "docProps/app.xml"}
    with zipfile.ZipFile(nowarn, "w", zipfile.ZIP_DEFLATED) as z:
        for n, data in items.items():
            z.writestr(n, data)
    ck("missing app.xml is WARN, tolerated without --strict", 0, run(nowarn, strict=False))

    print("----")
    print(f"test_pptx_mac_compat: {_pass} passed, {_fail} failed")
    return 0 if _fail == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
