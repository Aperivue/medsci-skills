#!/usr/bin/env python3
"""Build an editable STROBE participant flow diagram (.pptx) from a YAML config.

Why this script exists
----------------------
The Graphviz path (``generate_flow_diagram.R --type strobe``) renders an
auto-fitting monochrome diagram suitable for journal submission, but
co-authors often want to nudge box positions, edit prose, or recolor stage
labels in PowerPoint. This script produces a fully-editable .pptx in the
classical Identification → Screening → Inclusion → Analysis convention used
by the BMJ / Radiology / KJR cohort-study examples (Randolph 2018,
PLoS ONE 2021).

Unlike PRISMA, STROBE flow diagrams have study-specific spine structure
(number of cohort-stage boxes and exclusions varies per study), so this is
a single parametric builder rather than the PRISMA two-step
(``build_prisma2020_template.py`` then ``fill_prisma_template.py``).

YAML schema
-----------
    output_pptx: figures/figure1_strobe.pptx          # required
    slide_size: [13.33, 10.0]                          # inches; default widescreen
    title: "Figure 1. STROBE participant flow diagram"  # optional
    stages:
      - {name: Identification, color: "#1F3A68"}      # color = stage box fill
      - {name: Screening,      color: "#1F3A68"}
      - {name: Inclusion,      color: "#1F3A68"}
      - {name: Analysis,       color: "#1F3A68"}
    spine:
      - {id: enrolled,    stage: Identification, text: "..."}
      - {id: polyp_ever,  stage: Screening,      text: "..."}
      - {id: primary,     stage: Inclusion,      text: "..."}
      - {id: ksar,        stage: Analysis,       text: "..."}
      - {id: landmark,    stage: Analysis,       text: "..."}  # consecutive same-stage rows share one stage label
    exclusions:
      - {after: enrolled,   text: "Excluded (n = 147,245):\\n• no gallbladder polyp on any ultrasound"}
      - {after: polyp_ever, text: "Excluded (n = 3):\\n• prior C23 / zero post-baseline FU"}

Usage
-----
    python3 build_strobe_template.py \
        --config figures/figure1_strobe.yaml \
        --out    figures/figure1_strobe.pptx

Open the resulting .pptx in PowerPoint to fine-tune positions or styling
before saving as PDF / TIFF for submission.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


# ── Defaults ────────────────────────────────────────────────────────────────
DEFAULT_SLIDE = (13.33, 10.0)

# Column geometry (inches)
PHASE_X = 0.4
PHASE_W = 1.5
SPINE_X = 2.4
SPINE_W = 4.4
EXCL_X  = 7.4
EXCL_W  = 4.6

# Vertical spacing
TITLE_Y = 0.25
TITLE_H = 0.5
TOP_PAD = 0.35           # below title before first row
ROW_GAP = 0.3            # vertical gap between adjacent rows
DEFAULT_ROW_H = 1.05
DEFAULT_EXCL_H = 0.95

# Colors
NAVY   = RGBColor(0x1F, 0x3A, 0x68)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
LBLUE  = RGBColor(0xCF, 0xE1, 0xF5)


# ── Helpers ─────────────────────────────────────────────────────────────────
def _parse_color(hex_str: str | None, default: RGBColor) -> RGBColor:
    if not hex_str:
        return default
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _readable_text_color(bg: RGBColor) -> RGBColor:
    """Return BLACK or WHITE depending on background luminance."""
    # Relative luminance per WCAG (sRGB approx without gamma)
    r, g, b = bg[0], bg[1], bg[2]
    lum = 0.299 * r + 0.587 * g + 0.114 * b
    return RGBColor(0x1F, 0x3A, 0x68) if lum > 160 else WHITE


def add_box(slide, x, y, w, h, text, *,
            fill=WHITE, border=BLACK, font_color=BLACK,
            font_size=11, bold_first=True, anchor_middle=True,
            align_center=True, line_pt=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(line_pt)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    if anchor_middle:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align_center:
            para.alignment = 2  # center
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bool(bold_first and i == 0)
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=BLACK, width_pt=1.25):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    # Arrow head on the destination end
    line_xml = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(line_xml, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


def add_right_angle_arrow(slide, x1, y1, x2, y2, *, color=BLACK, width_pt=1.0):
    """Spine right edge → exclusion left edge with a horizontal stub then arrow."""
    # Use a single straight horizontal connector since both ends share y
    add_arrow(slide, x1, y1, x2, y2, color=color, width_pt=width_pt)


# ── Loader ──────────────────────────────────────────────────────────────────
def load_config(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    if path.suffix in (".yaml", ".yml"):
        if not HAS_YAML:
            sys.exit("PyYAML not installed; install or use JSON config.")
        return yaml.safe_load(text)
    if path.suffix == ".json":
        return json.loads(text)
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        if HAS_YAML:
            return yaml.safe_load(text)
        sys.exit("Config not JSON and PyYAML unavailable.")


# ── Builder ────────────────────────────────────────────────────────────────
def build(cfg: dict, out_path: Path) -> None:
    slide_w, slide_h = cfg.get("slide_size", DEFAULT_SLIDE)

    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_text = cfg.get("title")
    if title_text:
        title_box = slide.shapes.add_textbox(
            Inches(PHASE_X), Inches(TITLE_Y),
            Inches(slide_w - PHASE_X * 2), Inches(TITLE_H),
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = NAVY

    # Layout: compute y-coordinates for each spine row
    spine = cfg["spine"]
    n_spine = len(spine)
    spine_h = float(cfg.get("spine_box_height", DEFAULT_ROW_H))
    excl_h  = float(cfg.get("exclusion_box_height", DEFAULT_EXCL_H))

    # Available vertical space below title
    avail_top = TITLE_Y + TITLE_H + TOP_PAD if title_text else TITLE_Y
    row_pitch = spine_h + ROW_GAP
    # If too tall, scale down spine_h
    needed = avail_top + n_spine * row_pitch + 0.4
    if needed > slide_h:
        scale = (slide_h - avail_top - 0.4) / (n_spine * row_pitch)
        spine_h = max(0.7, spine_h * scale)
        excl_h  = max(0.6, excl_h * scale)
        row_pitch = spine_h + ROW_GAP

    spine_y = {}
    for i, b in enumerate(spine):
        spine_y[b["id"]] = avail_top + i * row_pitch

    # Phase column — group consecutive same-stage rows under one merged label
    stage_color_lookup = {s["name"]: _parse_color(s.get("color"), NAVY) for s in cfg.get("stages", [])}
    stages_seq = [b["stage"] for b in spine]
    i = 0
    while i < n_spine:
        j = i
        while j + 1 < n_spine and stages_seq[j + 1] == stages_seq[i]:
            j += 1
        sname = stages_seq[i]
        y_top = spine_y[spine[i]["id"]]
        y_bot = spine_y[spine[j]["id"]] + spine_h
        h = y_bot - y_top
        stage_fill = stage_color_lookup.get(sname, NAVY)
        add_box(
            slide, PHASE_X, y_top, PHASE_W, h,
            sname,
            fill=stage_fill,
            border=stage_fill,
            font_color=_readable_text_color(stage_fill),
            font_size=14,
            bold_first=True,
            line_pt=0.0,
        )
        i = j + 1

    # Spine boxes + arrows
    for k, b in enumerate(spine):
        y = spine_y[b["id"]]
        add_box(
            slide, SPINE_X, y, SPINE_W, spine_h,
            b["text"],
            fill=WHITE, border=BLACK, font_color=BLACK,
            font_size=11, bold_first=True,
        )
        if k > 0:
            prev_y = spine_y[spine[k - 1]["id"]]
            add_arrow(
                slide,
                SPINE_X + SPINE_W / 2, prev_y + spine_h,
                SPINE_X + SPINE_W / 2, y,
            )

    # Exclusion boxes + connector arrows
    spine_id_to_idx = {b["id"]: i for i, b in enumerate(spine)}
    for excl in cfg.get("exclusions", []):
        after_id = excl["after"]
        idx = spine_id_to_idx[after_id]
        y_after = spine_y[after_id]
        # Place exclusion at the same y as the "after" spine box (right-side branch)
        add_box(
            slide, EXCL_X, y_after, EXCL_W, excl_h,
            excl["text"],
            fill=WHITE, border=BLACK, font_color=BLACK,
            font_size=10, bold_first=False,
            align_center=False,  # left-align text in exclusions
        )
        # Connector: spine right edge → excl left edge (same y, mid-row)
        add_arrow(
            slide,
            SPINE_X + SPINE_W, y_after + spine_h / 2,
            EXCL_X,             y_after + excl_h / 2,
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


# ── CLI ─────────────────────────────────────────────────────────────────────
def main():
    p = argparse.ArgumentParser(description="Build an editable STROBE flow diagram .pptx from a YAML config.")
    p.add_argument("--config", required=True, type=Path, help="YAML/JSON STROBE config")
    p.add_argument("--out",    required=True, type=Path, help="Output .pptx path")
    args = p.parse_args()

    cfg = load_config(args.config)
    if not isinstance(cfg, dict):
        sys.exit(f"Config root must be a mapping; got {type(cfg)}")

    build(cfg, args.out)
    print(f"Wrote {args.out}")


if __name__ == "__main__":
    main()
