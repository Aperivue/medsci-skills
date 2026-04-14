#!/usr/bin/env python3
"""
MedSci Skills Demo 2: BCG Vaccine Meta-Analysis
Step 4 — Create Academic Presentation (present-paper skill)

Generates a 16:9 academic slide deck with embedded figures and speaker notes.

Usage: python3 04_create_pptx.py
"""

import sys
import logging
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).resolve().parent
OUTPUT = BASE
FIGURES = BASE / "figures"
LOGS = BASE


logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.StreamHandler(sys.stdout),
    ],
)
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x72, 0xB2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(1.1)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_header_bar(slide, title_text: str):
    """Add navy header bar with white title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, HEADER_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Teal accent line at bottom of header
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), HEADER_H - Inches(0.04),
        SLIDE_W, Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15), Inches(12), Inches(0.85)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.color.rgb = WHITE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Calibri"


def add_body_text(slide, text: str, left=0.6, top=1.4, width=12.0, height=5.5,
                  font_size=18, bold=False, color=DARK_TEXT, bullet=False):
    """Add body text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        line = line.strip()
        if not line:
            p.space_after = Pt(6)
            continue

        is_bullet = line.startswith("- ") or line.startswith("* ")
        if is_bullet:
            line = line[2:]
            p.level = 0

        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.font.bold = bold
        p.space_after = Pt(4)

    return txBox


def add_speaker_notes(slide, text: str):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_footer(slide, slide_num: int, total: int):
    """Add subtle footer with slide number."""
    txBox = slide.shapes.add_textbox(
        Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# Slide definitions
# ---------------------------------------------------------------------------
SLIDES = []

def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full navy background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Teal accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.3), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Efficacy of BCG Vaccination for\nPrevention of Tuberculosis"
    p.font.color.rgb = WHITE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "A Meta-Analysis of Randomized Controlled Trials"
    p2.font.color.rgb = TEAL
    p2.font.size = Pt(24)
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    # Author / date
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "MedSci Skills Demo"
    p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p3.font.size = Pt(18)
    p3.font.name = "Calibri"
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf3.add_paragraph()
    p4.text = f"Based on Colditz et al. (1994) | {date.today()}"
    p4.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p4.font.size = Pt(14)
    p4.font.name = "Calibri"
    p4.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide,
        "Welcome. This presentation summarizes a meta-analysis of 13 randomized "
        "controlled trials evaluating BCG vaccine efficacy for tuberculosis prevention. "
        "The analysis is based on the classic Colditz et al. 1994 dataset and serves "
        "as a teaching demonstration of meta-analytic methods."
    )
    return slide


def slide_background(prs):
    """Slide 2: Background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Background")
    add_body_text(slide, (
        "- TB remains a leading infectious cause of death (1.3 million deaths/year)\n"
        "- BCG vaccine: only licensed TB vaccine, given to >100 million children annually\n"
        "- Reported efficacy ranges from 0% to >80% across studies\n"
        "- Landmark meta-analysis by Colditz et al. (JAMA 1994) identified geographic latitude as key moderator\n"
        "- Hypothesis: environmental non-tuberculous mycobacteria (NTM) exposure at lower latitudes attenuates BCG efficacy\n"
        "- Need: systematic quantification with modern meta-analytic methods"
    ), top=1.5, font_size=20)
    add_speaker_notes(slide,
        "Tuberculosis kills over a million people annually. BCG is the only licensed vaccine "
        "but its efficacy is highly variable. Colditz and colleagues published a seminal "
        "meta-analysis in 1994 showing that latitude strongly predicts vaccine efficacy. "
        "The leading explanation involves cross-reactive immunity from environmental NTM "
        "being more prevalent near the equator."
    )
    return slide


def slide_objectives(prs):
    """Slide 3: Objectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Objectives")
    add_body_text(slide, (
        "1. Quantify the overall protective efficacy of BCG vaccination against tuberculosis\n\n"
        "2. Assess the degree and sources of between-study heterogeneity\n\n"
        "3. Evaluate the effect of geographic latitude on vaccine efficacy via meta-regression\n\n"
        "4. Test the robustness of findings through sensitivity and publication bias analyses"
    ), top=1.8, font_size=22)
    add_speaker_notes(slide,
        "Our four objectives cover the full meta-analytic pipeline: pooled estimation, "
        "heterogeneity assessment, moderator analysis via meta-regression, and robustness "
        "checks including leave-one-out sensitivity and publication bias tests."
    )
    return slide


def slide_methods(prs):
    """Slide 4: Methods."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Methods")

    # Left column
    add_body_text(slide, (
        "Data Source\n"
        "- 13 RCTs from Colditz et al. (1994)\n"
        "- metafor::dat.bcg dataset\n"
        "- Published 1948-1980\n\n"
        "Effect Measure\n"
        "- Risk Ratio (RR) from 2x2 tables\n"
        "- Log-transformed for analysis"
    ), left=0.6, top=1.5, width=5.5, font_size=17)

    # Right column
    add_body_text(slide, (
        "Statistical Analysis\n"
        "- Random-effects model (REML)\n"
        "- Heterogeneity: Q, I-squared, tau-squared\n"
        "- Subgroup: allocation method\n"
        "- Meta-regression: absolute latitude\n"
        "- Publication bias: Egger, Begg, trim-and-fill\n"
        "- Sensitivity: leave-one-out\n"
        "- Software: R + metafor package"
    ), left=6.8, top=1.5, width=5.8, font_size=17)

    add_speaker_notes(slide,
        "We used the well-characterized dat.bcg dataset containing 13 RCTs. "
        "Risk ratios were computed from 2x2 tables and analyzed with a REML "
        "random-effects model. Our analysis pipeline includes heterogeneity "
        "assessment, subgroup analysis by allocation method, meta-regression "
        "with absolute latitude, and comprehensive publication bias and "
        "sensitivity analyses."
    )
    return slide


def slide_forest(prs):
    """Slide 5: Forest Plot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Results: Forest Plot")

    # Embed forest plot image
    img_path = FIGURES / "forest_plot.png"
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.8)
        )
    else:
        add_body_text(slide, "[Forest plot image not found]", top=3.0, font_size=24,
                      color=RGBColor(0xCC, 0x00, 0x00))

    add_speaker_notes(slide,
        "The forest plot shows all 13 studies with their individual risk ratios and "
        "95% confidence intervals. The pooled RR is 0.489, meaning BCG vaccination "
        "reduced TB risk by approximately 51%. Note the substantial variability across "
        "studies, with I-squared of 92.2%. The diamond at the bottom represents the "
        "pooled estimate under the REML random-effects model. Several studies, "
        "particularly TPT Madras and Comstock 1969/1976, show risk ratios near or "
        "above 1.0, indicating no protective effect in those settings."
    )
    return slide


def slide_heterogeneity(prs):
    """Slide 6: Heterogeneity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Heterogeneity Assessment")

    # Key metrics in a structured layout
    add_body_text(slide, (
        "Overall Pooled Estimate\n"
        "- RR = 0.489 (95% CI: 0.344 - 0.696)\n"
        "- Prediction interval: 0.155 - 1.549\n\n"
        "Heterogeneity Metrics\n"
        "- Q = 152.23 (df = 12, p < 0.001)\n"
        "- I-squared = 92.2% (substantial heterogeneity)\n"
        "- Tau-squared = 0.3132\n\n"
        "Subgroup Analysis (by allocation method)\n"
        "- Random (k=7): RR = 0.379 (0.221 - 0.650), I-squared = 89.9%\n"
        "- Alternate (k=2): RR = 0.582 (0.335 - 1.011), I-squared = 82.0%\n"
        "- Systematic (k=4): RR = 0.654 (0.323 - 1.324), I-squared = 86.4%"
    ), top=1.5, font_size=18)

    add_speaker_notes(slide,
        "Heterogeneity is substantial, with I-squared of 92.2%, meaning over 90% of "
        "variability is due to true differences between studies rather than sampling error. "
        "The prediction interval includes 1.0, cautioning that a future study might not "
        "show benefit. Randomly allocated trials showed the strongest protective effect, "
        "which may reflect better methodological quality. However, high residual "
        "heterogeneity persists within all subgroups."
    )
    return slide


def slide_metaregression(prs):
    """Slide 7: Meta-regression (bubble plot)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Meta-Regression: Latitude Effect")

    # Embed bubble plot
    img_path = FIGURES / "bubble_plot.png"
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.8)
        )
    else:
        add_body_text(slide, "[Bubble plot image not found]", top=3.0, font_size=24,
                      color=RGBColor(0xCC, 0x00, 0x00))

    # Key results on the right
    add_body_text(slide, (
        "Key Findings\n\n"
        "Coefficient: -0.0291\n"
        "(SE = 0.0072, p < 0.001)\n\n"
        "R-squared = 75.6%\n\n"
        "Higher latitude\n"
        "= greater efficacy\n\n"
        "Consistent with NTM\n"
        "exposure hypothesis"
    ), left=8.8, top=1.5, width=4.0, font_size=17, color=DARK_TEXT)

    add_speaker_notes(slide,
        "The bubble plot demonstrates the strong negative relationship between latitude "
        "and log risk ratio. Each bubble is sized proportionally to study precision. "
        "The regression coefficient of -0.029 means that for each degree increase in "
        "latitude, the log RR decreases by 0.029, reflecting greater vaccine efficacy. "
        "Latitude alone explains 75.6% of between-study variance. This is consistent "
        "with the hypothesis that NTM exposure at lower latitudes provides partial "
        "cross-reactive immunity, reducing the incremental benefit of BCG."
    )
    return slide


def slide_pub_bias(prs):
    """Slide 8: Publication Bias."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Publication Bias Assessment")

    # Embed funnel plot (trim-and-fill version)
    img_path = FIGURES / "funnel_trimfill.png"
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(0.5), Inches(1.3), Inches(7.5), Inches(5.5)
        )

    # Results on the right
    add_body_text(slide, (
        "Statistical Tests\n\n"
        "Egger's test:\n"
        "p = 0.189 (NS)\n\n"
        "Begg's test:\n"
        "p = 0.952 (NS)\n\n"
        "Trim-and-fill:\n"
        "1 study imputed\n"
        "Adjusted RR = 0.518\n"
        "(95% CI: 0.365-0.736)\n\n"
        "No significant\n"
        "publication bias detected"
    ), left=8.5, top=1.5, width=4.3, font_size=17, color=DARK_TEXT)

    add_speaker_notes(slide,
        "Publication bias assessment is reassuring. The funnel plot shows reasonable "
        "symmetry. Both Egger's regression test and Begg's rank correlation test are "
        "non-significant. Trim-and-fill analysis imputed only one study, and the adjusted "
        "estimate of 0.518 is very close to the original 0.489, both remaining significant. "
        "However, with only 13 studies, these tests have limited statistical power."
    )
    return slide


def slide_sensitivity(prs):
    """Slide 9: Sensitivity Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Sensitivity Analysis: Leave-One-Out")

    add_body_text(slide, (
        "Leave-One-Out Results\n"
        "- All 13 leave-one-out pooled RR estimates remain statistically significant\n"
        "- Range of pooled RR: 0.452 to 0.533\n"
        "- No single study reverses the overall conclusion\n\n"
        "Influence Diagnostics\n"
        "- No studies with |externally standardized residual| > 2\n"
        "- No unduly influential individual studies detected\n\n"
        "Interpretation\n"
        "- The protective effect of BCG vaccination is robust\n"
        "- Pooled estimate is not driven by any single trial\n"
        "- Results are stable across all sensitivity checks"
    ), top=1.5, font_size=19)

    add_speaker_notes(slide,
        "Leave-one-out analysis is a critical robustness check. By iteratively removing "
        "each study and re-fitting the model, we confirm that no single study unduly "
        "influences the pooled result. All 13 estimates remain significant, ranging from "
        "0.452 to 0.533. No studies had large externally standardized residuals. This "
        "gives us confidence that the overall conclusion of BCG efficacy is robust."
    )
    return slide


def slide_conclusions(prs):
    """Slide 10: Conclusions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Conclusions")

    add_body_text(slide, (
        "1. BCG vaccination significantly reduces TB risk by ~51%\n"
        "   (RR = 0.489, 95% CI: 0.344-0.696)\n\n"
        "2. Substantial heterogeneity (I-squared = 92.2%) driven primarily\n"
        "   by geographic latitude (R-squared = 75.6%)\n\n"
        "3. Greater efficacy at higher latitudes, consistent with NTM\n"
        "   exposure hypothesis\n\n"
        "4. No significant publication bias; results robust to\n"
        "   leave-one-out sensitivity analysis\n\n"
        "5. Implications: continued BCG use justified, especially at higher\n"
        "   latitudes; novel vaccines needed for tropical regions"
    ), top=1.5, font_size=20)

    add_speaker_notes(slide,
        "In summary, BCG vaccination provides significant protection against TB with "
        "about 51% risk reduction. However, efficacy varies dramatically by latitude, "
        "with this single moderator explaining three-quarters of the heterogeneity. "
        "The findings support continued BCG use, particularly in temperate and "
        "high-latitude regions, while highlighting the urgent need for next-generation "
        "vaccines effective in tropical settings where NTM exposure attenuates BCG benefit."
    )
    return slide


def slide_limitations(prs):
    """Slide 11: Limitations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Limitations")

    add_body_text(slide, (
        "- Historical dataset: trials conducted 1948-1980, evolving diagnostics and BCG strains\n\n"
        "- Aggregate data only: no individual patient data for within-study effect modification\n\n"
        "- Small number of studies (k=13): limits power for meta-regression and subgroup analyses\n\n"
        "- Teaching demonstration: uses curated dataset, not a de novo systematic review\n\n"
        "- No formal risk of bias assessment (RoB 2) or GRADE certainty evaluation\n\n"
        "- Wide prediction interval (0.155-1.549) includes 1.0, indicating uncertainty for new settings"
    ), top=1.5, font_size=19)

    add_speaker_notes(slide,
        "We should acknowledge several limitations. The data are historical, spanning "
        "a period when TB diagnostics and BCG manufacturing varied. We only have "
        "aggregate data, not individual patient data. With 13 studies, meta-regression "
        "power is limited. This is a re-analysis for teaching purposes, not a de novo "
        "systematic review. We also did not perform formal risk of bias or GRADE "
        "assessments, which would be required for a submission-ready manuscript."
    )
    return slide


def slide_thankyou(prs):
    """Slide 12: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full navy background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Teal accent
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.8), Inches(5.3), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Thank you text
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.3), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.color.rgb = WHITE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(4.2), Inches(9.3), Inches(2.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Questions?"
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.font.size = Pt(24)
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf2.add_paragraph()
    p3.text = f"\nDataset: metafor::dat.bcg | Reference: Colditz et al. (JAMA 1994)"
    p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p3.font.size = Pt(14)
    p3.font.name = "Calibri"
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf2.add_paragraph()
    p4.text = f"Generated with MedSci Skills | {date.today()}"
    p4.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
    p4.font.size = Pt(12)
    p4.font.name = "Calibri"
    p4.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide,
        "Thank you for your attention. I am happy to take any questions. "
        "The full analysis code and outputs are available as part of the "
        "MedSci Skills demo package."
    )
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    log.info("=" * 60)
    log.info("Step 4: Create Academic Presentation")
    log.info("=" * 60)

    prs = Presentation()

    # Set 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build slides
    slide_builders = [
        ("Title", slide_title),
        ("Background", slide_background),
        ("Objectives", slide_objectives),
        ("Methods", slide_methods),
        ("Forest Plot", slide_forest),
        ("Heterogeneity", slide_heterogeneity),
        ("Meta-regression", slide_metaregression),
        ("Publication Bias", slide_pub_bias),
        ("Sensitivity", slide_sensitivity),
        ("Conclusions", slide_conclusions),
        ("Limitations", slide_limitations),
        ("Thank You", slide_thankyou),
    ]

    total = len(slide_builders)
    for i, (name, builder) in enumerate(slide_builders, 1):
        log.info(f"  Building slide {i}/{total}: {name}")
        s = builder(prs)
        # Add footer (skip title and thank-you slides)
        if i not in (1, total):
            add_footer(s, i, total)

    # Save
    out_path = OUTPUT / "presentation.pptx"
    prs.save(str(out_path))
    log.info(f"Saved: {out_path}")
    log.info(f"  Slides: {total}")
    log.info(f"  File size: {out_path.stat().st_size / 1024:.1f} KB")
    log.info("=" * 60)
    log.info("Presentation generation complete.")
    log.info("=" * 60)


if __name__ == "__main__":
    main()
