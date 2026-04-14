"""
MedSci Skills Demo 3: NHANES Obesity & Diabetes
Step 5 — Academic Presentation (present-paper skill)

Generates a 11-slide 16:9 academic presentation with embedded figures.

Usage: python3 05_create_pptx.py
Output: output/presentation.pptx
"""

import os
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

print("=" * 60)
print("MedSci Skills Demo 3: Presentation Generation")
print(f"Date: {datetime.date.today()}")
print("=" * 60)

# === Constants ===
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x72, 0xB2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank


def add_navy_header(slide, title_text):
    """Add a navy header bar across the top of the slide."""
    # Header background
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"

    # Teal accent line
    line = slide.shapes.add_shape(
        1, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_body_text(slide, text, left=0.8, top=1.6, width=11.5, height=5.0,
                  font_size=20, color=DARK_GRAY, bold=False):
    """Add body text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.font.bold = bold
        p.space_after = Pt(8)
    return txBox


def add_bullet_slide(slide, title, bullets, notes_text=""):
    """Create a standard bullet-point slide."""
    add_navy_header(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6),
                                      Inches(11.5), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(12)
        p.level = 0
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def add_footer(slide):
    """Add a subtle footer."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.0),
                                      Inches(11.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "NHANES 2017-2018 | Obesity & Diabetes Analysis"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


# =========================================================
# SLIDE 1: Title Slide
# =========================================================
slide = prs.slides.add_slide(blank_layout)

# Full navy background
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Teal accent bar
accent = slide.shapes.add_shape(1, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06))
accent.fill.solid()
accent.fill.fore_color.rgb = TEAL
accent.line.fill.background()

# Title
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Association Between Body Mass Index\nand Diabetes Mellitus in US Adults"
p.font.size = Pt(36)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.0), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "A Cross-Sectional Analysis of NHANES 2017-2018"
p2.font.size = Pt(24)
p2.font.color.rgb = TEAL
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER

# Source
txBox3 = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.0), Inches(1.0))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "Data source: CDC National Health and Nutrition Examination Survey"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
p3.font.name = "Calibri"
p3.alignment = PP_ALIGN.CENTER

slide.notes_slide.notes_text_frame.text = (
    "Welcome to this presentation on the association between BMI and diabetes "
    "in US adults. This analysis uses data from NHANES 2017-2018, a nationally "
    "representative survey conducted by the CDC. Our goal is to quantify the "
    "dose-response relationship between BMI categories and diabetes prevalence "
    "using survey-weighted methods."
)

# =========================================================
# SLIDE 2: Background
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_bullet_slide(
    slide, "Background",
    [
        "Diabetes affects 537 million adults globally (IDF 2021)",
        "37.3 million Americans have diabetes; 90-95% are type 2",
        "Obesity is the most important modifiable risk factor for T2DM",
        "NHANES provides nationally representative US health data",
        "Complex survey design requires weighted analysis for valid estimates",
        "Continued surveillance is essential for tracking disparities",
    ],
    "Diabetes is a global health crisis with substantial morbidity, mortality, and economic "
    "burden. Obesity is the single most important modifiable risk factor. NHANES provides "
    "the gold-standard data for estimating disease burden in the US, but the complex "
    "survey design must be properly accounted for to produce valid estimates."
)
add_footer(slide)

# =========================================================
# SLIDE 3: Objectives
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_bullet_slide(
    slide, "Study Objectives",
    [
        "Primary: Quantify the association between BMI categories and diabetes",
        "  prevalence in US adults using NHANES 2017-2018",
        "",
        "Secondary objectives:",
        "  1. Estimate survey-weighted diabetes prevalence by BMI category",
        "  2. Calculate adjusted odds ratios controlling for demographics",
        "  3. Examine racial/ethnic disparities in the BMI-diabetes relationship",
        "  4. Demonstrate the impact of survey weights on prevalence estimates",
    ],
    "Our primary objective is to quantify the dose-response relationship between "
    "BMI and diabetes. Secondary objectives include estimating weighted prevalence, "
    "computing adjusted odds ratios, examining disparities, and demonstrating why "
    "survey weights matter."
)
add_footer(slide)

# =========================================================
# SLIDE 4: Methods
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_bullet_slide(
    slide, "Methods",
    [
        "Design: Cross-sectional analysis of NHANES 2017-2018",
        "Population: US adults aged >= 20 years (n = 4,866)",
        "Exposure: BMI (WHO categories: Normal, Overweight, Obese)",
        "Outcome: Diabetes defined as HbA1c >= 6.5% (ADA criteria)",
        "Covariates: Age, sex, race/ethnicity, education level",
        "Analysis: Survey-weighted logistic regression (WTMEC2YR)",
        "  - Model 1: Unadjusted (BMI category only)",
        "  - Model 2: Adjusted for all covariates",
    ],
    "This is a cross-sectional study using one cycle of NHANES data. We included "
    "4,866 adults with complete BMI and HbA1c data. BMI was categorized per WHO "
    "definitions. Diabetes was defined using the ADA HbA1c threshold of 6.5%. "
    "We used survey-weighted logistic regression with two models: unadjusted and "
    "adjusted for age, sex, race/ethnicity, and education."
)
add_footer(slide)

# =========================================================
# SLIDE 5: Table 1 Highlights
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_navy_header(slide, "Study Population (Table 1)")

# Create a simple text-based table summary
table_data = [
    ["", "Normal\n(n=1,189)", "Overweight\n(n=1,593)", "Obese\n(n=2,084)", "p-value"],
    ["Age (years)", "49.5 +/- 19.3", "53.5 +/- 17.1", "51.2 +/- 16.7", "<0.001"],
    ["Female, %", "54.4%", "46.3%", "54.8%", "<0.001"],
    ["BMI (kg/m2)", "22.4 +/- 1.7", "27.4 +/- 1.4", "36.5 +/- 6.3", "<0.001"],
    ["HbA1c (%)", "5.60 +/- 0.88", "5.85 +/- 1.12", "6.05 +/- 1.17", "<0.001"],
    ["Diabetes, %", "7.5%", "13.9%", "19.9%", "<0.001"],
]

n_rows = len(table_data)
n_cols = len(table_data[0])
tbl = slide.shapes.add_table(n_rows, n_cols,
                              Inches(0.8), Inches(1.5),
                              Inches(11.5), Inches(4.0)).table

for r in range(n_rows):
    for c in range(n_cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK_GRAY
                if c == 0:
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.font.bold = True

    # Header row coloring
    if r == 0:
        for c in range(n_cols):
            tbl.cell(r, c).fill.solid()
            tbl.cell(r, c).fill.fore_color.rgb = NAVY

# Key takeaway
add_body_text(slide,
    "Key: 42.8% of participants were obese; diabetes prevalence nearly tripled from normal to obese",
    left=0.8, top=6.0, width=11.5, height=0.6, font_size=16, color=TEAL, bold=True)

slide.notes_slide.notes_text_frame.text = (
    "Table 1 shows the baseline characteristics. Nearly 43% of participants were obese, "
    "reflecting the high obesity prevalence in the US. Mean age was about 51 years with "
    "roughly equal sex distribution. Notice the clear gradient in HbA1c and diabetes "
    "prevalence across BMI categories. All differences were statistically significant "
    "with p < 0.001."
)
add_footer(slide)

# =========================================================
# SLIDE 6: Prevalence Results (Figure 1)
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_navy_header(slide, "Diabetes Prevalence by BMI Category")

fig_path = "figures/prevalence_by_bmi.png"
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(1.5), Inches(1.5), Inches(5.5))

# Key numbers on the right
add_body_text(slide,
    "Survey-Weighted Prevalence:\n\n"
    "Normal:       4.1% (6.1-9.1%)\n"
    "Overweight:  8.8% (12.3-15.7%)\n"
    "Obese:        14.7% (18.2-21.6%)\n\n"
    "Overall:       10.2% (13.9-15.9%)\n\n"
    "Unweighted overall: 14.9%\n"
    "(4.7 pp higher without weights)",
    left=7.5, top=1.5, width=5.0, height=5.0, font_size=18)

slide.notes_slide.notes_text_frame.text = (
    "Figure 1 shows diabetes prevalence by BMI category. There is a clear dose-response "
    "pattern: prevalence roughly doubles from normal to overweight and again from overweight "
    "to obese. The survey-weighted prevalence is notably lower than unweighted, reflecting "
    "the oversampling of high-risk groups in NHANES. This 4.7 percentage point difference "
    "demonstrates why survey weights are essential."
)
add_footer(slide)

# =========================================================
# SLIDE 7: Regression Results (Figure 2)
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_navy_header(slide, "Adjusted Odds Ratios for Diabetes")

fig_path = "figures/or_forest_plot.png"
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(0.5), Inches(1.5), Inches(7.0))

# Key ORs on the right
add_body_text(slide,
    "Key Adjusted ORs:\n\n"
    "Obese vs Normal:\n  OR 4.50 (4.49-4.51)\n\n"
    "Overweight vs Normal:\n  OR 2.06 (2.05-2.06)\n\n"
    "Age (per year):\n  OR 1.06 (1.06-1.06)\n\n"
    "Female vs Male:\n  OR 0.70 (0.70-0.70)\n\n"
    "All p < 0.001",
    left=8.0, top=1.5, width=4.5, height=5.0, font_size=17)

slide.notes_slide.notes_text_frame.text = (
    "The forest plot shows adjusted odds ratios from the multivariable model. "
    "Obesity was associated with 4.5-fold higher odds of diabetes, and overweight "
    "with 2-fold higher odds. Age increased risk by 6% per year. Females had 30% "
    "lower odds than males. All associations were statistically significant."
)
add_footer(slide)

# =========================================================
# SLIDE 8: Subgroup Analysis (Figure 4)
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_navy_header(slide, "Prevalence by Age Group and BMI")

fig_path = "figures/prevalence_by_age_bmi.png"
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(1.5), Inches(1.5), Inches(5.5))

add_body_text(slide,
    "Key Observations:\n\n"
    "- Diabetes prevalence increases\n  with both age and BMI\n\n"
    "- Even in young adults (20-39),\n  obesity elevates risk\n\n"
    "- Strongest absolute difference\n  in the 60-79 age group\n\n"
    "- Consistent dose-response\n  pattern across all age strata",
    left=7.5, top=1.5, width=5.0, height=5.0, font_size=18)

slide.notes_slide.notes_text_frame.text = (
    "This figure shows that the BMI-diabetes gradient is consistent across age groups. "
    "Importantly, even among younger adults aged 20-39, obesity is associated with "
    "substantially higher diabetes prevalence. The absolute difference is largest in "
    "older adults, where the combined effects of aging and obesity are most pronounced."
)
add_footer(slide)

# =========================================================
# SLIDE 9: Racial/Ethnic Disparities
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_bullet_slide(
    slide, "Racial/Ethnic Disparities",
    [
        "Adjusted ORs compared with Non-Hispanic White:",
        "",
        "  Non-Hispanic Asian:        OR 2.97 (2.96-2.97)",
        "  Non-Hispanic Black:        OR 1.84 (1.83-1.84)",
        "  Mexican American:           OR 1.58 (1.57-1.58)",
        "",
        "NH Asian: highest adjusted risk despite lowest mean BMI",
        "  - Supports ethnicity-specific BMI thresholds (WHO 2004)",
        "  - Metabolic risk at lower BMI in Asian populations",
        "",
        "Disparities persist after adjusting for BMI, age, sex, education",
    ],
    "Racial and ethnic disparities in diabetes risk persisted even after adjustment for BMI "
    "and demographic factors. Non-Hispanic Asians had the highest adjusted odds ratio, which "
    "is particularly notable given their lower mean BMI. This finding aligns with evidence "
    "that Asian populations develop metabolic complications at lower BMI thresholds, supporting "
    "calls for ethnicity-specific cutoffs. Non-Hispanic Black and Mexican American populations "
    "also showed significantly elevated risk."
)
add_footer(slide)

# =========================================================
# SLIDE 10: Discussion
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_bullet_slide(
    slide, "Discussion",
    [
        "Findings consistent with dose-response BMI-diabetes relationship",
        "Adjusted OR 4.50 for obesity comparable to prior NHANES analyses",
        "",
        "Methodological insight: survey weights are critical",
        "  - Unweighted: 14.9% vs Weighted: 10.2% (4.7 pp difference)",
        "  - NHANES oversamples minorities and older adults",
        "",
        "Limitations:",
        "  - Cross-sectional design (no causal inference)",
        "  - HbA1c-only diabetes definition (may underestimate)",
        "  - Cannot distinguish type 1 from type 2 diabetes",
        "  - Residual confounding (diet, physical activity, family history)",
    ],
    "Our findings are consistent with decades of epidemiological evidence. The effect "
    "size for obesity is comparable to prior estimates. An important methodological "
    "lesson is the 4.7 percentage point difference between weighted and unweighted "
    "prevalence, demonstrating why survey weights must be used. Key limitations include "
    "the cross-sectional design, the HbA1c-only definition, and potential residual confounding."
)
add_footer(slide)

# =========================================================
# SLIDE 11: Conclusions
# =========================================================
slide = prs.slides.add_slide(blank_layout)
add_bullet_slide(
    slide, "Conclusions",
    [
        "Clear dose-response relationship: BMI category -> diabetes prevalence",
        "",
        "Obesity: 4.5x higher adjusted odds of diabetes vs normal weight",
        "Overweight: 2.0x higher adjusted odds vs normal weight",
        "",
        "Significant racial/ethnic disparities persist after adjustment",
        "  - Asian populations: highest risk at lowest BMI",
        "",
        "Survey weights are essential for NHANES analyses",
        "",
        "Public health implication: obesity prevention remains the",
        "  cornerstone strategy for reducing the diabetes burden",
    ],
    "In conclusion, we demonstrated a strong, graded association between BMI and diabetes "
    "in a nationally representative US sample. Obesity conferred 4.5-fold higher odds of "
    "diabetes. Racial and ethnic disparities persisted after comprehensive adjustment. "
    "These findings reinforce the critical importance of obesity prevention and weight "
    "management as key strategies for reducing the population burden of diabetes."
)
add_footer(slide)

# =========================================================
# SLIDE 12: Thank You
# =========================================================
slide = prs.slides.add_slide(blank_layout)

# Full navy background
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Teal accent
accent = slide.shapes.add_shape(1, Inches(4.0), Inches(3.5), Inches(5.333), Inches(0.04))
accent.fill.solid()
accent.fill.fore_color.rgb = TEAL
accent.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.0), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(44)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.0), Inches(2.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Data: NHANES 2017-2018 (CDC/NCHS)"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER

p3 = tf2.add_paragraph()
p3.text = "Analysis pipeline: Python 3 | statsmodels | matplotlib"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
p3.font.name = "Calibri"
p3.alignment = PP_ALIGN.CENTER

p4 = tf2.add_paragraph()
p4.text = "Generated with MedSci Skills for Claude Code"
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
p4.font.name = "Calibri"
p4.alignment = PP_ALIGN.CENTER

slide.notes_slide.notes_text_frame.text = (
    "Thank you for your attention. I am happy to take any questions about the "
    "methodology, findings, or implications of this analysis."
)

# === Save ===

prs.save("presentation.pptx")
print(f"Saved: presentation.pptx ({len(prs.slides)} slides)")
print("=" * 60)
